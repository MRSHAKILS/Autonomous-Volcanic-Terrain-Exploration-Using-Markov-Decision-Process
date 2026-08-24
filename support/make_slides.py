"""Generate the project presentation PPTX files (final and update decks).

The course manual requires every slide to carry a footer with the group
number and the names of all group members, and requires the presentation to
open by introducing all members by full name. This script builds both decks
programmatically so they can be regenerated at any time:

    python support/make_slides.py

It writes ``others/final_presentation.pptx`` and
``others/update_presentation.pptx``, embedding figures from ``images/``.
"""

import sys
from pathlib import Path


if __package__ is None or __package__ == "":
    sys.path.append(str(Path(__file__).resolve().parent.parent))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parent.parent
IMAGES = PROJECT_ROOT / "images"
OTHERS = PROJECT_ROOT / "others"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

DARK = RGBColor(0x15, 0x18, 0x21)
GOLD = RGBColor(0xE8, 0xC1, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x9F, 0xB2, 0xC8)
INK = RGBColor(0x20, 0x24, 0x2B)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)

GROUP = "Group 5"
MEMBERS = "Shakil Ahmed  ·  Fahim Foysal  ·  Shefa Tabassum  ·  Tanvir Ahmed"
FOOTER_TEXT = f"{GROUP}  |  {MEMBERS}"


def _new_deck() -> Presentation:
    deck = Presentation()
    deck.slide_width = SLIDE_WIDTH
    deck.slide_height = SLIDE_HEIGHT
    return deck


def _blank_slide(deck: Presentation, dark: bool = False):
    slide = deck.slides.add_slide(deck.slide_layouts[6])  # blank layout

    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK if dark else WHITE

    _add_footer(slide, dark=dark)
    return slide


def _add_footer(slide, dark: bool) -> None:
    """Mandatory footer: group number and the names of all group members."""
    box = slide.shapes.add_textbox(Inches(0.3), SLIDE_HEIGHT - Inches(0.42), SLIDE_WIDTH - Inches(0.6), Inches(0.35))
    paragraph = box.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = FOOTER_TEXT
    run.font.size = Pt(11)
    run.font.color.rgb = GRAY if dark else RGBColor(0x60, 0x66, 0x70)


def _add_text(slide, left, top, width, height, lines, dark=False):
    """Add a textbox from (text, size_pt, bold, color_or_None) tuples."""
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True

    for index, (text, size, bold, color) in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(10)
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color if color is not None else (WHITE if dark else INK)

    return box


def _title_line(slide, text, dark=False):
    _add_text(slide, Inches(0.55), Inches(0.35), SLIDE_WIDTH - Inches(1.1), Inches(0.9), [(text, 30, True, None)], dark=dark)


def _bullets(slide, items, left=Inches(0.65), top=Inches(1.45), width=Inches(6.6), size=17):
    _add_text(slide, left, top, width, Inches(5.4), [(f"•  {item}", size, False, None) for item in items])


def _picture(slide, image_name, left, top, width):
    slide.shapes.add_picture(str(IMAGES / image_name), left, top, width=width)


def _title_slide(deck, subtitle_lines):
    slide = _blank_slide(deck, dark=True)
    _add_text(
        slide,
        Inches(0.8),
        Inches(1.7),
        SLIDE_WIDTH - Inches(1.6),
        Inches(4.5),
        [
            ("Autonomous Volcanic Terrain Exploration", 38, True, WHITE),
            ("Using a Markov Decision Process (MDP)", 26, False, GOLD),
            ("", 12, False, GRAY),
        ]
        + subtitle_lines
        + [
            ("", 12, False, GRAY),
            ("We are Group 5:", 18, True, WHITE),
            ("Shakil Ahmed   ·   Fahim Foysal   ·   Shefa Tabassum   ·   Tanvir Ahmed", 18, False, GRAY),
            ("CSE 440 (Artificial Intelligence)  ·  Section 1  ·  Instructor: Dr. Mohammad Shifat-E-Rabbi", 13, False, GRAY),
        ],
        dark=True,
    )
    return slide


# ---------------------------------------------------------------------------
# Final presentation.
# ---------------------------------------------------------------------------


def build_final_deck() -> Path:
    deck = _new_deck()

    _title_slide(deck, [("Final Project Presentation", 18, False, WHITE)])

    slide = _blank_slide(deck)
    _title_line(slide, "The problem")
    _bullets(slide, [
        "Volcanic terrain is too dangerous for humans to explore directly",
        "An autonomous rover must collect science samples and survive",
        "Lava is fatal; craters and gas are costly; rock blocks movement",
        "Movement is unreliable: the rover can slip sideways or stall",
        "The terrain itself changes during the mission",
    ])
    _picture(slide, "hazards_before.png", Inches(7.5), Inches(1.3), Inches(5.3))

    slide = _blank_slide(deck)
    _title_line(slide, "The task as an MDP")
    _bullets(slide, [
        "States: every walkable cell of the grid (rock excluded)",
        "Actions: up / down / left / right / stay / scan",
        "Transitions: 75% intended, 10% drift each side, 5% stall",
        "Rewards: science +20, base +5, safe −1, gas −15, crater −40, lava −100",
        "Discount γ = 0.9, tuned so distant science never justifies deadly gambles",
        "Solved exactly with value iteration (converges in milliseconds)",
    ])
    _add_text(slide, Inches(7.6), Inches(2.2), Inches(5.2), Inches(2.5), [
        ("V(s) ← max over a of  Σ P(s′|s,a) [ R(s,a,s′) + γ V(s′) ]", 16, False, ACCENT),
        ("The Bellman backup, repeated until convergence", 13, False, None),
    ])

    slide = _blank_slide(deck)
    _title_line(slide, "The computed policy")
    _bullets(slide, [
        "One optimal action per cell: a complete plan, not a single path",
        "The flow field bends around lava and craters",
        "It converges on science points and the base",
        "Robust by construction: slips land the agent on cells that already have a plan",
    ], width=Inches(6.0))
    _picture(slide, "policy_map.png", Inches(6.9), Inches(1.15), Inches(6.0))

    slide = _blank_slide(deck)
    _title_line(slide, "Collect-once science + re-planning")
    _bullets(slide, [
        "Naive MDP behavior: camp forever on one +20 science cell",
        "Our rule: a sample is collected once, then the cell becomes ordinary ground",
        "The agent immediately re-runs value iteration (warm start: a few sweeps)",
        "Result: an efficient, risk-adjusted tour of the objectives, then return to base",
    ])

    slide = _blank_slide(deck)
    _title_line(slide, "Dynamic hazards")
    _bullets(slide, [
        "Gas clouds drift into neighboring safe cells (p = 0.15 per cloud per step)",
        "Lava spills into safe neighbors (p = 0.08) and cools back after 6 steps",
        "Base, science, rock, and craters are never destroyed, so the state space stays fixed",
        "Every terrain change triggers a fresh re-plan",
        "A lava flow reaching the rover destroys it",
    ], width=Inches(6.0))
    _picture(slide, "hazards_before.png", Inches(6.8), Inches(1.5), Inches(3.15))
    _picture(slide, "hazards_after.png", Inches(10.05), Inches(1.5), Inches(3.15))
    _add_text(slide, Inches(6.8), Inches(5.3), Inches(6.3), Inches(0.5), [
        ("Same map: initial (left) and after 20 hazard steps, 58 change events (right)", 13, False, None),
    ])

    slide = _blank_slide(deck)
    _title_line(slide, "A completed mission")
    _bullets(slide, [
        "Seed 42, dynamic hazards ON",
        "24 steps, 2 science samples, survived",
        "71 terrain-change events during the mission",
        "Reward higher than the static run of the same seed: re-planning exploited a gas cloud drifting off its route",
    ], width=Inches(5.6))
    _picture(slide, "final_map.png", Inches(6.6), Inches(1.1), Inches(6.3))

    slide = _blank_slide(deck)
    _title_line(slide, "Does planning beat heuristics? 20 random terrains")
    _bullets(slide, [
        "MDP: follows the value-iteration policy",
        "Greedy: BFS to nearest unvisited cell, ignoring soft hazards",
        "Random: uniform random direction",
        "Only the MDP agent earns positive reward: +42 vs −124 vs −180",
        "It survives 60% of missions (2× greedy) with ~5× fewer hazard entries",
    ], width=Inches(5.5), size=16)
    _picture(slide, "performance_plot.png", Inches(6.4), Inches(1.1), Inches(6.5))

    slide = _blank_slide(deck)
    _title_line(slide, "Conclusions & future work")
    _bullets(slide, [
        "Decision-theoretic planning decisively beats heuristics in hazardous exploration",
        "Safety and science behavior emerge from the reward model, with no hand-coded rules",
        "The reward model IS the specification: 'camping' was optimal until we changed the rules",
        "Future: risk-sensitive objectives, energy budgets, fog-of-war POMDP, multi-rover teams",
        "Everything is seeded and reproducible: python main.py",
    ])

    output_path = OTHERS / "final_presentation.pptx"
    deck.save(str(output_path))
    return output_path


# ---------------------------------------------------------------------------
# Update presentation (Update 1 state of the project, 12th class).
# ---------------------------------------------------------------------------


def build_update_deck() -> Path:
    deck = _new_deck()

    _title_slide(deck, [("Project Update 1", 18, False, WHITE)])

    slide = _blank_slide(deck)
    _title_line(slide, "Problem & goal")
    _bullets(slide, [
        "An autonomous rover exploring hazardous volcanic terrain",
        "Grid world: safe ground, lava, craters, gas, rock, science points, base",
        "Movement is stochastic: the rover can slip or stall",
        "Goal: collect science, avoid hazards, and survive, planned with an MDP",
    ])
    _picture(slide, "terrain_map_seed_1.png", Inches(7.6), Inches(1.5), Inches(5.2))

    slide = _blank_slide(deck)
    _title_line(slide, "Done: terrain generator & MDP core (Member 1)")
    _bullets(slide, [
        "Seeded procedural terrain with configurable cell probabilities",
        "Hazard density capped so maps stay traversable; CSV save/load",
        "MDP formulation: states, 6 actions, slip transitions, reward model",
        "Value iteration + optimal policy extraction, γ = 0.9",
    ], width=Inches(6.2))
    _picture(slide, "terrain_map_seed_7.png", Inches(7.2), Inches(1.4), Inches(5.4))

    slide = _blank_slide(deck)
    _title_line(slide, "Done: agent, simulation & visualization (Members 2 & 3)")
    _bullets(slide, [
        "MDPExplorerAgent: starts at base, follows the policy, samples real outcomes",
        "Tracks reward, hazards, science, survival; mission summary",
        "Simulation loop with step budget and coverage target",
        "Matplotlib terrain rendering with legend, one color per cell type",
    ])

    slide = _blank_slide(deck)
    _title_line(slide, "Done: baseline experiments (Member 4)")
    _bullets(slide, [
        "Three agents compared on identical terrains: MDP vs Greedy vs Random",
        "Shared metrics: reward, coverage, hazards, science, survival",
        "CSV results + performance plot generated automatically",
    ], width=Inches(5.6), size=16)
    _picture(slide, "performance_plot.png", Inches(6.4), Inches(1.2), Inches(6.5))

    slide = _blank_slide(deck)
    _title_line(slide, "Plan to the final submission")
    _bullets(slide, [
        "Dynamic hazards: drifting gas and spreading lava, with re-planning",
        "Fix the science-camping behavior (collect-once rule in the core agent)",
        "Mission map with the real agent path instead of a sample path",
        "Full pipeline in main.py; final report, slides, one-minute demo video",
    ])

    output_path = OTHERS / "update_presentation.pptx"
    deck.save(str(output_path))
    return output_path


def main() -> None:
    print(f"Saved: {build_final_deck()}")
    print(f"Saved: {build_update_deck()}")


if __name__ == "__main__":
    main()
