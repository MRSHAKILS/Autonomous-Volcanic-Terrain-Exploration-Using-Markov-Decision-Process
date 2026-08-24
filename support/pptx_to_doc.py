"""Convert the project presentation PPTX files to Word documents.

Reads the generated decks in ``others/`` and rebuilds each one as a .docx:
every slide becomes a section with its heading, its bullet points, and its
embedded images, so the same content is available to anyone who prefers Word
over PowerPoint. Run it after ``make_slides.py``:

    python support/pptx_to_doc.py

It writes ``others/final_presentation.docx`` and
``others/update_presentation.docx``.
"""

import io
import sys
from pathlib import Path


if __package__ is None or __package__ == "":
    sys.path.append(str(Path(__file__).resolve().parent.parent))

from docx import Document
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


PROJECT_ROOT = Path(__file__).resolve().parent.parent
OTHERS = PROJECT_ROOT / "others"

DECKS = ["final_presentation.pptx", "update_presentation.pptx"]

GROUP_LINE = "Group 5  |  Shakil Ahmed  ·  Fahim Foysal  ·  Shefa Tabassum  ·  Tanvir Ahmed"
HEADING_SIZE_THRESHOLD_PT = 24  # slide text at or above this size becomes a doc heading
IMAGE_WIDTH = Inches(4.6)


def _shape_paragraphs(shape) -> list[tuple[str, float]]:
    """Return (text, font_size_pt) for each non-empty paragraph of a shape."""
    if not shape.has_text_frame:
        return []

    paragraphs = []
    for paragraph in shape.text_frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs).strip()
        if not text:
            continue

        sizes = [run.font.size.pt for run in paragraph.runs if run.font.size is not None]
        paragraphs.append((text, max(sizes) if sizes else 12.0))

    return paragraphs


def _add_slide_to_document(document: Document, slide, slide_number: int) -> None:
    document.add_heading(f"Slide {slide_number}", level=1)

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            document.add_picture(io.BytesIO(shape.image.blob), width=IMAGE_WIDTH)
            continue

        for text, size in _shape_paragraphs(shape):
            if text == GROUP_LINE:
                continue  # the per-slide footer is carried once in the doc header

            if text.startswith("•"):
                document.add_paragraph(text.lstrip("• ").strip(), style="List Bullet")
            elif size >= HEADING_SIZE_THRESHOLD_PT:
                document.add_heading(text, level=2)
            else:
                document.add_paragraph(text)


def convert(deck_name: str) -> Path:
    """Convert one PPTX deck in others/ to a Word document alongside it."""
    deck_path = OTHERS / deck_name
    presentation = Presentation(str(deck_path))

    document = Document()

    # Carry the mandatory group/member identification in the page header.
    header_paragraph = document.sections[0].header.paragraphs[0]
    header_paragraph.text = GROUP_LINE
    for run in header_paragraph.runs:
        run.font.size = Pt(9)

    for slide_number, slide in enumerate(presentation.slides, start=1):
        _add_slide_to_document(document, slide, slide_number)

    output_path = deck_path.with_suffix(".docx")
    document.save(str(output_path))
    return output_path


def main() -> None:
    for deck_name in DECKS:
        print(f"Saved: {convert(deck_name)}")


if __name__ == "__main__":
    main()
